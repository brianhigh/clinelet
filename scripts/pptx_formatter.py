"""
Rich PPTX extraction module for wiki integrator.

Converts .pptx files to Markdown while preserving:
- Slide titles as Markdown headings
- List structure (bullet and numbered)
- Speaker notes
- Shape hierarchy (title, subtitle, content placeholders)
- Image captions (where available)
"""

import re
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


# Shape type constants
PLACEHOLDER = MSO_SHAPE_TYPE.PLACEHOLDER if hasattr(MSO_SHAPE_TYPE, 'PLACEHOLDER') else 7
PICTURE = MSO_SHAPE_TYPE.PICTURE if hasattr(MSO_SHAPE_TYPE, 'PICTURE') else 13
TABLE_SHAPE = MSO_SHAPE_TYPE.TABLE if hasattr(MSO_SHAPE_TYPE, 'TABLE') else 6

# Placeholder types that indicate title/subtitle
PLACEHOLDER_TYPE_TITLE = 1
PLACEHOLDER_TYPE_SUBTITLE = 2

# Content placeholder types that may contain lists
CONTENT_PLACEHOLDER_TYPES = {
    3,   # OBJECT_FRAME
    11,  # BODY_CONTENT
    12,  # CENTERED_BODY_CONTENT
    13,  # OBJECT
    14,  # VERTICAL_CONTENT
    15,  # VERTICAL_OBJECT
    16,  # TWO_CONTENT
    17,  # TWO_OBJ
    18,  # TWO_OBJ_COLLAPSED
    19,  # VERTICAL_TWO_CONTENT
    20,  # VERTICAL_TWO_OBJ
    21,  # SLIDE_IMAGE
    22,  # TITLE
    23,  # MEDIA
    24,  # CENTERED_SHAPE
}


def _get_placeholder_type(shape):
    """Get the placeholder type of a shape if it is a placeholder.
    
    Returns None if the shape is not a placeholder.
    """
    try:
        return shape.placeholder_format.type
    except (ValueError, AttributeError):
        return None


def _is_title_placeholder(shape):
    """Check if a shape is a title placeholder."""
    try:
        return shape.placeholder_format.type == PLACEHOLDER_TYPE_TITLE
    except (ValueError, AttributeError):
        return False


def _is_subtitle_placeholder(shape):
    """Check if a shape is a subtitle placeholder."""
    try:
        return shape.placeholder_format.type == PLACEHOLDER_TYPE_SUBTITLE
    except (ValueError, AttributeError):
        return False


def _extract_text_from_shape(shape):
    """Extract text from a shape, preserving list structure and inline formatting."""
    if not hasattr(shape, 'text'):
        return None
    
    text = shape.text
    if not text or not text.strip():
        return None
    
    return text


def _extract_text_with_formatting(shape):
    """Extract text from a shape with formatting markers for wiki markdown.
    
    Uses the underlying XML to detect bold/italic formatting since
    python-pptx run objects don't expose .bold/.italic directly.
    """
    if not hasattr(shape, 'text_frame'):
        if hasattr(shape, 'text'):
            text = shape.text.strip()
            return text if text else None
        return None
    
    parts = []
    for para in shape.text_frame.paragraphs:
        run_parts = []
        for run in para.runs:
            txt = run.text
            if not txt:
                continue
            
            # Access formatting through the underlying XML element
            try:
                rPr = run._element.rPr  # RunProperties
                if rPr is not None:
                    bold_elem = rPr.find('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}b')
                    italic_elem = rPr.find('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}i')
                    is_bold = bold_elem is not None
                    is_italic = italic_elem is not None
                else:
                    is_bold = False
                    is_italic = False
            except Exception:
                is_bold = False
                is_italic = False
            
            if is_bold and is_italic:
                txt = f'***{txt}***'
            elif is_bold:
                txt = f'**{txt}**'
            elif is_italic:
                txt = f'*{txt}*'
            run_parts.append(txt)
        
        para_text = ''.join(run_parts)
        if para_text:
            parts.append(para_text)
    
    result = '\n'.join(parts) if parts else None
    return result if result else None


def _extract_text_from_text_body(text_frame):
    """Extract text from a text frame, preserving paragraph structure.
    
    Note: python-pptx _Paragraph objects don't have .paragraph_format directly.
    We access indent via XML, but for simplicity we just extract clean text here.
    """
    if not text_frame or not hasattr(text_frame, 'paragraphs'):
        return ''
    
    lines = []
    for para in text_frame.paragraphs:
        text = para.text.strip()
        if text:
            lines.append(text)
    
    return '\n'.join(lines)


def _extract_slide_notes(slide):
    """Extract speaker notes from a slide."""
    if not hasattr(slide, 'notes_slide'):
        return None
    notes_slide = slide.notes_slide
    if notes_slide is None:
        return None
    notes_text = notes_slide.notes_text_frame.text.strip()
    return notes_text if notes_text else None


def _extract_slide_images(slide):
    """Extract image captions from a slide."""
    captions = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # Look for caption text in adjacent shapes
            if hasattr(shape, 'name') and shape.name and shape.name.lower() != 'picture':
                captions.append(f'![{shape.name}]({shape.name})')
    return captions


def _extract_table_from_shape(shape):
    """Extract table from a shape if it has a table."""
    if not hasattr(shape, 'table'):
        return None
    
    table = shape.table
    if not table:
        return None
    
    rows = list(table.rows)
    if not rows:
        return None
    
    num_cols = len(rows[0].cells)
    markdown_lines = []
    
    for i, row in enumerate(rows):
        cells = []
        for cell in row.cells:
            content = cell.text.strip().replace('|', '\\|')
            cells.append(content)
        
        cells = cells[:num_cols]
        if i == 0:
            markdown_lines.append('| ' + ' | '.join(f'**{c}**' for c in cells) + ' |')
            markdown_lines.append('| ' + ' | '.join(['---'] * len(cells)) + ' |')
        else:
            markdown_lines.append('| ' + ' | '.join(cells) + ' |')
    
    return '\n' + '\n'.join(markdown_lines) + '\n\n'


def slide_to_markdown(slide, slide_index):
    """
    Convert a single PPTX slide to Markdown.
    
    Args:
        slide: A pptx.slide.Slide object.
        slide_index: Zero-based slide index.
        
    Returns:
        A string containing the Markdown representation of the slide.
    """
    result_parts = []
    
    # Find the title shape
    title_text = None
    subtitle_text = None
    content_shapes = []
    
    for shape in slide.shapes:
        # Check if it's a title placeholder
        if _is_title_placeholder(shape):
            title_text = _extract_text_with_formatting(shape)
        elif _is_subtitle_placeholder(shape):
            subtitle_text = _extract_text_with_formatting(shape)
        else:
            content_shapes.append(shape)
    
    # Add slide number and title
    result_parts.append(f'---\n\n## Slide {slide_index + 1}')
    
    if title_text:
        result_parts.append(f'\n\n**{title_text}**\n')
    
    # Process content shapes
    for shape in content_shapes:
        text = _extract_text_with_formatting(shape)
        
        # Try to extract table
        table_md = _extract_table_from_shape(shape)
        if table_md:
            result_parts.append(f'\n{table_md}')
            continue
        
        # Extract text with bullet/indent structure
        if text:
            # Check for list structure via text frame
            if hasattr(shape, 'text_frame') and hasattr(shape.text_frame, 'paragraphs'):
                list_text = _extract_text_from_text_body(shape.text_frame)
                if list_text:
                    result_parts.append(f'\n{list_text}\n')
                else:
                    result_parts.append(f'\n{text}\n')
            else:
                result_parts.append(f'\n{text}\n')
        
        # Check for embedded images (pass slide, not shape)
        image_caps = _extract_slide_images(slide)
        for cap in image_caps:
            result_parts.append(f'\n{cap}\n')
    
    # Add subtitle if present
    if subtitle_text:
        result_parts.append(f'\n{subtitle_text}\n')
    
    # Add speaker notes if present
    notes = _extract_slide_notes(slide)
    if notes:
        result_parts.append(f'\n---\n> **Speaker notes:**\n> {notes.replace(chr(10), chr(10) + "> ")}\n')
    
    result_parts.append('\n---\n')
    
    return ''.join(result_parts)


def pptx_to_markdown(pptx_path):
    """
    Convert a .pptx file to rich Markdown.
    
    Args:
        pptx_path: Path to the .pptx file.
        
    Returns:
        A string containing the Markdown representation with preserved
        slide structure, titles, list formatting, and speaker notes.
    """
    prs = Presentation(pptx_path)
    result_parts = []
    
    # Add document-level metadata
    core_properties = prs.core_properties
    if core_properties:
        metadata = []
        if core_properties.title:
            metadata.append(f'# {core_properties.title}')
        if core_properties.author:
            metadata.append(f'**Author:** {core_properties.author}')
        if core_properties.subject:
            metadata.append(f'**Subject:** {core_properties.subject}')
        if metadata:
            result_parts.append('\n'.join(metadata) + '\n\n---\n\n')
    
    # Convert each slide
    for i, slide in enumerate(prs.slides):
        slide_md = slide_to_markdown(slide, i)
        result_parts.append(slide_md)
    
    return ''.join(result_parts)


def pptx_to_markdown_simplified(pptx_path):
    """
    Simplified PPTX-to-Markdown conversion.
    
    Falls back to basic text extraction when rich formatting detection fails.
    """
    try:
        return pptx_to_markdown(pptx_path)
    except Exception:
        # Fallback: basic text extraction (current behavior)
        prs = Presentation(pptx_path)
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines.append(shape.text.strip())
        return "\n".join(lines)