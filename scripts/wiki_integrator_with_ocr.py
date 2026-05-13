import os
import re
import sys
import subprocess
import tempfile
import shutil
import platform
import struct
import zipfile

# Constants
RAW_DIR = "raw"
SPACE_DIR = "wiki"
STATE_DIR = ".clinelet"
MANIFEST_PATH = os.path.join(STATE_DIR, "processed_files.txt")

# Dependency tracking
MISSING_DEPENDENCIES = []

# Magic bytes (file signatures) for MIME type detection
# Format: {offset_bytes: (signature_bytes, mime_type)}
MAGIC_SIGNATURES = {
    # PDF
    b'%PDF': 'application/pdf',
    # DOCX (ZIP-based, magic at offset 0x50)
    b'PK\x03\x04': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    # XLSX (ZIP-based)
    b'PK\x03\x04': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
    # PPTX (ZIP-based)
    b'PK\x03\x04': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
    # PNG
    b'\x89PNG\r\n\x1a\n': 'image/png',
    # JPEG
    b'\xff\xd8\xff': 'image/jpeg',
    # GIF
    b'GIF87a': 'image/gif',
    b'GIF89a': 'image/gif',
    # BMP
    b'BM': 'image/bmp',
    # TIFF (little-endian)
    b'II\x2a\x00': 'image/tiff',
    # TIFF (big-endian)
    b'MM\x00\x2a': 'image/tiff',
    # WebP
    b'RIFF': 'image/webp',  # RIFF header, MIME confirmed by checking for 'WEBP' at offset 8
    # HTML (also covers XML)
    b'<!DOCTYPE html': 'text/html',
    b'<html': 'text/html',
    b'<?xml': 'text/xml',
    # ZIP (general container for DOCX/XLSX/PPTX)
    b'PK\x03\x04': 'application/zip',
}

# MIME type to extension mapping for validation
MIME_TO_EXTENSIONS = {
    'application/pdf': ['.pdf'],
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': ['.docx'],
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': ['.xlsx'],
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': ['.pptx'],
    'image/png': ['.png'],
    'image/jpeg': ['.jpg', '.jpeg'],
    'image/gif': ['.gif'],
    'image/bmp': ['.bmp'],
    'image/tiff': ['.tiff', '.tif'],
    'image/webp': ['.webp'],
    'text/plain': ['.txt', '.md'],
    'text/html': ['.html', '.htm'],
    'text/xml': ['.xml'],
    'application/zip': ['.zip'],
    'application/rtf': ['.rtf'],
    'text/rtf': ['.rtf'],
}

# Try to import python-magic for more robust MIME detection
HAVE_PYTHON_MAGIC = False
try:
    import magic
    HAVE_PYTHON_MAGIC = True
except ImportError:
    pass

def get_os_info():
    """Detect the operating system and return appropriate package manager info.
    
    Returns a dict with:
        - os_type: 'linux', 'darwin', 'windows'
        - primary_cmd: Primary install command template
        - alternatives: List of (cmd, label) tuples for alternative managers
    """
    system = platform.system()
    
    if system == "Linux":
        # Detect which package managers are available
        alternatives = []
        primary_cmd = None
        
        # Check for apt (Debian/Ubuntu)
        if shutil.which("apt"):
            primary_cmd = "sudo apt install {package}"
        # Check for dnf (Fedora/RHEL)
        elif shutil.which("dnf"):
            primary_cmd = "sudo dnf install {package}"
        # Check for yum (older Fedora/RHEL)
        elif shutil.which("yum"):
            primary_cmd = "sudo yum install {package}"
        # Check for pacman (Arch)
        elif shutil.which("pacman"):
            primary_cmd = "sudo pacman -S {package}"
        # Check for zypper (openSUSE)
        elif shutil.which("zypper"):
            primary_cmd = "sudo zypper install {package}"
        
        # Add alternatives if we have a primary
        if primary_cmd:
            if "apt" in primary_cmd:
                alternatives = [
                    ("sudo dnf install {package}", "Fedora/RHEL"),
                    ("sudo yum install {package}", "Older Fedora/RHEL"),
                    ("sudo pacman -S {package}", "Arch Linux"),
                    ("sudo zypper install {package}", "openSUSE"),
                ]
            elif "dnf" in primary_cmd:
                alternatives = [
                    ("sudo apt install {package}", "Debian/Ubuntu"),
                    ("sudo pacman -S {package}", "Arch Linux"),
                ]
            elif "pacman" in primary_cmd:
                alternatives = [
                    ("sudo apt install {package}", "Debian/Ubuntu"),
                    ("sudo dnf install {package}", "Fedora/RHEL"),
                ]
        
        if primary_cmd is None:
            primary_cmd = "sudo <package-manager> install {package}  # Detect your package manager (apt/dnf/pacman/zypper)"
            alternatives = [
                ("sudo apt install {package}", "Debian/Ubuntu"),
                ("sudo dnf install {package}", "Fedora/RHEL"),
                ("sudo pacman -S {package}", "Arch Linux"),
                ("sudo zypper install {package}", "openSUSE"),
            ]
        
        return {
            "os_type": "linux",
            "primary_cmd": primary_cmd,
            "alternatives": alternatives,
        }
    
    elif system == "Darwin":
        # macOS
        alternatives = []
        primary_cmd = None
        
        # Check for brew (Homebrew)
        if shutil.which("brew"):
            primary_cmd = "brew install {package}"
            alternatives = [
                ("sudo port install {package}", "MacPorts"),
            ]
        # Check for port (MacPorts)
        elif shutil.which("port"):
            primary_cmd = "sudo port install {package}"
            alternatives = [
                ("brew install {package}", "Homebrew"),
            ]
        
        if primary_cmd is None:
            primary_cmd = "brew install {package}  # Install Homebrew: https://brew.sh"
            alternatives = [
                ("sudo port install {package}", "MacPorts"),
            ]
        
        return {
            "os_type": "darwin",
            "primary_cmd": primary_cmd,
            "alternatives": alternatives,
        }
    
    elif system == "Windows":
        # Windows
        alternatives = []
        primary_cmd = None
        
        # Check for winget
        if shutil.which("winget"):
            primary_cmd = "winget install {package}"
            alternatives = [
                ("choco install {package}", "Chocolatey"),
                ("scoop install {package}", "Scoop"),
            ]
        # Check for choco (Chocolatey)
        elif shutil.which("choco"):
            primary_cmd = "choco install {package}"
            alternatives = [
                ("winget install {package}", "winget"),
                ("scoop install {package}", "Scoop"),
            ]
        # Check for scoop
        elif shutil.which("scoop"):
            primary_cmd = "scoop install {package}"
            alternatives = [
                ("winget install {package}", "winget"),
                ("choco install {package}", "Chocolatey"),
            ]
        
        if primary_cmd is None:
            primary_cmd = "winget install {package}  # Install winget: https://microsoft.github.io/winget"
            alternatives = [
                ("choco install {package}", "Chocolatey"),
                ("scoop install {package}", "Scoop"),
            ]
        
        return {
            "os_type": "windows",
            "primary_cmd": primary_cmd,
            "alternatives": alternatives,
        }
    
    else:
        return {
            "os_type": "unknown",
            "primary_cmd": "{package_manager} install {package}",
            "alternatives": [],
        }

def load_manifest():
    """Loads processed filenames and ensures state directory exists."""
    if not os.path.exists(STATE_DIR):
        os.makedirs(STATE_DIR)
        
    if not os.path.exists(MANIFEST_PATH):
        # Create the file if it doesn't exist
        open(MANIFEST_PATH, 'a').close()
        return set()

    try:
        with open(MANIFEST_PATH, 'r', encoding='utf-8') as f:
            return set(line.strip() for line in f if line.strip())
    except Exception as e:
        print(f"Error reading manifest: {e}")
        return set()

def save_to_manifest(filename):
    """Appends a filename to the manifest file."""
    try:
        if not os.path.exists(STATE_DIR):
            os.makedirs(STATE_DIR)
            
        with open(MANIFEST_PATH, 'a', encoding='utf-8') as f:
            f.write(filename + "\n")
    except Exception as e:
        print(f"Warning: Could not update manifest: {e}")

def detect_file_type(file_path):
    """Detect the MIME type of a file by examining its magic bytes (file signature).
    
    This provides content-based detection rather than relying on filename extensions,
    which helps catch mismatched or misleading file types.
    
    Args:
        file_path: Path to the file to detect.
        
    Returns:
        A string containing the detected MIME type (e.g., 'application/pdf'),
        or None if the type could not be determined.
    """
    # If python-magic is available, use it for best results
    if HAVE_PYTHON_MAGIC:
        try:
            mime = magic.Magic(mime=True)
            detected = mime.from_file(file_path)
            if detected:
                return detected
        except Exception:
            pass
    
    # Fallback to manual magic byte detection
    try:
        with open(file_path, 'rb') as f:
            # Read enough bytes to cover all our signatures
            header = f.read(12)
            if not header or len(header) < 2:
                return None
            
            # Check fixed signatures first (exact byte match at offset 0)
            for sig, mime_type in MAGIC_SIGNATURES.items():
                if isinstance(sig, bytes) and header[:len(sig)] == sig:
                    # Special case: RIFF -> could be WebP, need to check further
                    if sig == b'RIFF':
                        if len(header) >= 12 and f.read(4) == b'WEBP':
                            return 'image/webp'
                        else:
                            # RIFF without WEBP - might be audio, skip
                            continue
                    return mime_type
            
            # For ZIP-based files (DOCX, XLSX, PPTX), we need to inspect the ZIP contents
            if header[:4] == b'PK\x03\x04':
                try:
                    with zipfile.ZipFile(file_path, 'r') as zf:
                        for name in zf.namelist():
                            name_lower = name.lower()
                            if name_lower.startswith('word/') and name_lower.endswith('.xml'):
                                return 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
                            elif name_lower.startswith('xl/') and name_lower.endswith('.xml'):
                                return 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
                            elif name_lower.startswith('ppt/') and name_lower.endswith('.xml'):
                                return 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
                        # Generic ZIP
                        return 'application/zip'
                except Exception:
                    return 'application/zip'
            
            # For text files, check the first line more carefully
            try:
                first_line = header.decode('utf-8', errors='ignore').strip().lower()
                if first_line.startswith('<!doctype html') or first_line.startswith('<html'):
                    return 'text/html'
                elif first_line.startswith('<?xml'):
                    return 'text/xml'
                elif first_line and all(c.isprintable() or c in '\n\r\t' for c in first_line[:100]):
                    # Looks like plain text
                    return 'text/plain'
            except Exception:
                pass
            
            return None
    except Exception:
        return None


def validate_mime_vs_extension(file_path, detected_mime):
    """Validate that the file's extension matches its detected MIME type.
    
    Args:
        file_path: Path to the file.
        detected_mime: The detected MIME type string.
        
    Returns:
        A tuple (is_valid, warning_message). is_valid is True if the extension
        matches the MIME type. warning_message describes any mismatch.
    """
    ext = os.path.splitext(file_path)[1].lower()
    
    if detected_mime not in MIME_TO_EXTENSIONS:
        # Unknown MIME type, skip validation
        return True, None
    
    expected_extensions = MIME_TO_EXTENSIONS[detected_mime]
    
    if ext in expected_extensions:
        return True, None
    else:
        # Check if the MIME type is a generic container (ZIP) - allow flexibility
        if detected_mime == 'application/zip':
            return True, None
        if detected_mime == 'text/plain':
            # Plain text is flexible - .md, .txt, .html, .xml all count
            return True, None
        
        return False, f"Extension '{ext}' does not match detected MIME type '{detected_mime}'. Expected one of: {expected_extensions}"


def to_snake_case(filename):
    """Converts filename to snake_case and ensures it ends in .md"""
    name = os.path.splitext(filename)[0]
    name = re.sub(r'[^a-zA-Z0-9]+', '_', name)
    return name.strip('_').lower() + ".md"

def get_unique_path(target_dir, filename):
    """Prevents overwriting by appending a counter if the file exists."""
    base, ext = os.path.splitext(filename)
    counter = 1
    target_path = os.path.join(target_dir, filename)
    
    while os.path.exists(target_path):
        target_path = os.path.join(target_dir, f"{base}_{counter}{ext}")
        counter += 1
    return target_path

def _find_tool(name):
    """Find a tool by searching PATH first, then common installation directories.
    Returns the full path to the tool or None if not found.
    Designed to be portable across different systems.
    """
    # First try PATH
    result = shutil.which(name)
    if result:
        return result
    
    # Common Windows installation directories to search
    search_dirs = [
        r"C:\Program Files",
        r"C:\Program Files (x86)",
        r"C:\ProgramData",
        os.environ.get("LOCALAPPDIR", ""),
        os.environ.get("PROGRAMFILES", ""),
        os.environ.get("PROGRAMFILES(X86)", ""),
        os.environ.get("PROGRAMW6432", ""),
    ]
    
    # Filter out empty entries
    search_dirs = [d for d in search_dirs if d and os.path.isdir(d)]
    
    # Tool-specific search patterns
    tool_patterns = {
        "tesseract": [
            r"Tesseract-OCR\tesseract.exe",
        ],
        "pdftoppm": [
            r"poppler\Library\bin\pdftoppm.exe",
            r"poppler\bin\pdftoppm.exe",
            r"poppler-*/bin\pdftoppm.exe",
            r"Poppler\Library\bin\pdftoppm.exe",
            r"Poppler\bin\pdftoppm.exe",
        ],
        "pdftocairo": [
            r"poppler\Library\bin\pdftocairo.exe",
            r"poppler\bin\pdftocairo.exe",
            r"poppler-*/bin\pdftocairo.exe",
            r"Poppler\Library\bin\pdftocairo.exe",
            r"Poppler\bin\pdftocairo.exe",
        ],
    }
    
    patterns = tool_patterns.get(name, [f"{name}.exe" if os.name == 'nt' else name])
    
    for search_dir in search_dirs:
        for pattern in patterns:
            # Handle wildcard patterns like poppler-*
            if '*' in pattern:
                # Get the base directory and wildcard part
                base_subdir = pattern.split('/')[0] if '/' in pattern else pattern
                full_search = os.path.join(search_dir, base_subdir)
                if os.path.isdir(full_search):
                    for item in os.listdir(full_search):
                        candidate = os.path.join(search_dir, item, *pattern.split('/')[1:])
                        if os.path.exists(candidate):
                            return candidate
            else:
                candidate = os.path.join(search_dir, pattern)
                if os.path.exists(candidate):
                    return candidate
    
    return None

def get_tesseract_data_dir(tesseract_path):
    """Find the tessdata directory for a given tesseract installation.
    Returns the path to the tessdata directory or None.
    """
    tesseract_dir = os.path.dirname(tesseract_path)
    tessdata = os.path.join(tesseract_dir, "tessdata")
    if os.path.isdir(tessdata) and os.path.exists(os.path.join(tessdata, "eng.traineddata")):
        return tessdata
    
    # Also search parent directories
    parent = os.path.dirname(tesseract_dir)
    tessdata = os.path.join(parent, "tessdata")
    if os.path.isdir(tessdata) and os.path.exists(os.path.join(tessdata, "eng.traineddata")):
        return tessdata
    
    return None

def set_tesseract_env(tesseract_path):
    """Set TESSDATA_PREFIX environment variable for tesseract.
    Returns True if successful, False otherwise.
    """
    tessdata_dir = get_tesseract_data_dir(tesseract_path)
    if tessdata_dir:
        os.environ["TESSDATA_PREFIX"] = tessdata_dir
        return True
    print("[!] Warning: Could not find tessdata directory for tesseract")
    return False

def ocr_image(file_path):
    """OCR on an image file using tesseract, with pre-conversion to 8-bit RGB to avoid Tesseract/Leptonica issues."""
    tesseract_path = _find_tool("tesseract")
    if tesseract_path is None:
        print("[!] tesseract not found in PATH")
        return None
    
    magick_path = shutil.which("magick")
    
    with tempfile.TemporaryDirectory() as tmpdir:
        converted_img_path = os.path.join(tmpdir, "converted_image.png")
        conversion_success = False

        # Try Pillow first
        try:
            from PIL import Image
            with Image.open(file_path) as img:
                rgb_img = img.convert("RGB")
                rgb_img.save(converted_img_path, "PNG")
            conversion_success = True
        except Exception as e:
            print(f"[*] Pillow conversion failed for {file_path}, trying ImageMagick: {e}")
            
            # Fallback to ImageMagick
            if magick_path:
                try:
                    result = subprocess.run(
                        [magick_path, file_path, converted_img_path],
                        capture_output=True,
                        text=True,
                        check=True
                    )
                    conversion_success = True
                except subprocess.CalledProcessError as err:
                    print(f"[!] ImageMagick conversion failed for {file_path}: {err.stderr}")
            else:
                print("[!] ImageMagick (magick) not found in PATH. Cannot fallback.")

        if not conversion_success:
            print(f"[!] Failed to convert image {file_path} for OCR using both Pillow and ImageMagick.")
            return None

        base_path = os.path.join(tmpdir, "ocr_output")
        result = subprocess.run(
            [tesseract_path, converted_img_path, base_path],
            capture_output=True,
            text=True
        )
        if result.returncode != 0:
            return None
        txt_path = base_path + ".txt"
        if os.path.exists(txt_path):
            with open(txt_path, 'r', encoding='utf-8') as f:
                return f.read().strip()
    return None


def process_file(file_path, filename):
    """Extracts text and saves to markdown. Returns False if skipped/failed."""
    new_filename = to_snake_case(filename)
    target_path = get_unique_path(SPACE_DIR, new_filename)
    
    ext = os.path.splitext(filename)[1].lower()
    content = ""
    
    try:
        if ext in [".md", ".txt"]:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                content = f.read()

        elif ext == ".html":
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                html = f.read()
            content = re.sub(r'<[^>]+>', '', html)

        elif ext == ".pdf":
            try:
                import pypdf
                reader = pypdf.PdfReader(file_path)
                content = "\n".join([page.extract_text() or "" for page in reader.pages]).strip()
                
                if not content:
                    print(f"[*] PDF '{filename}' has no text layer. Attempting OCR...")
                    try:
                        pdftoppm_path = _find_tool("pdftoppm")
                        if pdftoppm_path is None:
                            print("[!] pdftoppm not found in PATH. Install poppler for PDF OCR support.")
                            return False

                        with tempfile.TemporaryDirectory() as tmpdir:
                            prefix = os.path.join(tmpdir, "page")
                            subprocess.run(
                                [pdftoppm_path, "-png", file_path, prefix],
                                check=True,
                                capture_output=True,
                                text=True
                            )
                            
                            # Get all png files and sort them numerically by the page number in filename
                            image_files = []
                            for f in os.listdir(tmpdir):
                                match = re.search(r'page-(\d+)\.png$', f)
                                if match:
                                    image_files.append((int(match.group(1)), f))
                            
                            image_files.sort()
                            
                            ocr_contents = []
                            for _, img_name in image_files:
                                img_path = os.path.join(tmpdir, img_name)
                                ocr_text = ocr_image(img_path)
                                ocr_contents.append(ocr_text if ocr_text else "")
                            
                            content = "\n".join(ocr_contents).strip()
                            
                            if not content:
                                print(f"[!] OCR failed to extract text from {filename}")
                                return False
                            else:
                                print(f"[+] OCR successful for {filename}")
                    except (subprocess.CalledProcessError, FileNotFoundError, Exception) as e:
                        print(f"[-] OCR failed for {filename}: {e}")
                        return False
            except ImportError:
                print(f"[-] Missing 'pypdf'. Skipping: {filename}")
                return False

        elif ext in [".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".tif", ".webp"]:
            ocr_text = ocr_image(file_path)
            if ocr_text:
                content = ocr_text
                print(f"[+] OCR successful for {filename}")
            else:
                print(f"[!] OCR failed to extract text from {filename}")
                return False

        elif ext == ".docx":
            try:
                import docx
                doc = docx.Document(file_path)
                content = "\n".join([p.text for p in doc.paragraphs]).strip()
            except ImportError:
                print(f"[-] Missing 'python-docx'. Skipping: {filename}")
                return False

        elif ext == ".xlsx":
            try:
                import openpyxl
                wb = openpyxl.load_workbook(file_path, data_only=True)
                lines = []
                for sheet in wb.worksheets:
                    lines.append(f"--- Sheet: {sheet.title} ---")
                    for row in sheet.iter_rows(values_only=True):
                        row_str = " ".join(str(c) for c in row if c is not None).strip()
                        if row_str: lines.append(row_str)
                content = "\n".join(lines)
            except ImportError:
                print(f"[-] Missing 'openpyxl'. Skipping: {filename}")
                return False

        elif ext == ".pptx":
            try:
                import pptx
                prs = pptx.Presentation(file_path)
                lines = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            lines.append(shape.text.strip())
                content = "\n".join(lines)
            except ImportError:
                print(f"[-] Missing 'python-pptx'. Skipping: {filename}")
                return False

        else:
            print(f"[!] Unsupported format: {ext} ({filename})")
            return False

        # Final check: If content is still empty after processing
        if not content.strip():
            print(f"[!] Skipped: {filename} (No content extracted)")
            return False

        # Final write to markdown
        with open(target_path, 'w', encoding='utf-8') as f:
            f.write(content)
        
        return True

    except Exception as e:
        print(f"[ERROR] Critical error processing {filename}: {e}")
        return False

def check_python_module(module_name, import_name=None, install_name=None):
    """Check if a Python module is available. Returns True if found.
    
    Args:
        module_name: The module import name (e.g., 'PIL')
        import_name: The actual import name if different from module_name
        install_name: The pip install name if different from module_name
    
    Adds to MISSING_DEPENDENCIES if not found.
    """
    if import_name is None:
        import_name = module_name
    if install_name is None:
        install_name = module_name
    
    try:
        __import__(import_name)
        return True
    except ImportError:
        MISSING_DEPENDENCIES.append({
            "type": "python",
            "module": module_name,
            "install": install_name,
            "description": f"Python module '{module_name}'"
        })
        return False


def check_system_tool(tool_name, package_name=None):
    """Check if a system tool is available. Returns True if found.
    
    Args:
        tool_name: The executable name (e.g., 'tesseract')
        package_name: The system package name for installation instructions.
                      If None, uses tool_name as the package name.
    
    Adds to MISSING_DEPENDENCIES if not found. Uses OS-specific install commands.
    """
    if shutil.which(tool_name) or _find_tool(tool_name):
        return True
    
    if package_name is None:
        package_name = tool_name
    
    # Get OS-specific install commands
    os_info = get_os_info()
    
    # Build the install command using the primary package manager
    primary_cmd = os_info["primary_cmd"].format(package=package_name)
    
    # Store alternatives for the user
    alt_commands = [cmd.format(package=package_name) for cmd, label in os_info["alternatives"]]
    
    MISSING_DEPENDENCIES.append({
        "type": "system",
        "tool": tool_name,
        "package": package_name,
        "install_cmd": primary_cmd,
        "alternatives": alt_commands,
        "description": f"System tool '{tool_name}' (package: {package_name})"
    })
    return False


def print_dependency_report():
    """Print a report of all missing dependencies with installation instructions."""
    if not MISSING_DEPENDENCIES:
        print("[✓] All dependencies are installed.")
        return False
    
    os_info = get_os_info()
    
    print("\n" + "=" * 60)
    print("[!] MISSING DEPENDENCIES DETECTED")
    print("=" * 60)
    print(f"\nDetected OS: {os_info['os_type'].title()}")
    print("The following dependencies are missing. Install them to enable")
    print("full functionality of the wiki integrator.\n")
    
    # Group by type
    python_deps = [d for d in MISSING_DEPENDENCIES if d["type"] == "python"]
    system_deps = [d for d in MISSING_DEPENDENCIES if d["type"] == "system"]
    
    if python_deps:
        print("Python Modules:")
        print("-" * 40)
        for dep in python_deps:
            print(f"  - {dep['description']}:")
            print(f"    Install: pip install {dep['install']}")
        print()
    
    if system_deps:
        print("System Packages:")
        print("-" * 40)
        for dep in system_deps:
            print(f"  - {dep['description']}:")
            print(f"    Install: {dep['install_cmd']}")
            # Provide alternative install commands from OS detection
            if dep.get("alternatives"):
                print("    Alternatives:")
                for alt in dep["alternatives"]:
                    print(f"      - {alt}")
        print()
    
    print("Quick install all Python modules:")
    if python_deps:
        pip_packages = " ".join(set(d["install"] for d in python_deps))
        print(f"  pip install {pip_packages}")
    print()
    
    if system_deps:
        packages = " ".join(set(d["package"] for d in system_deps))
        print(f"Quick install all system packages:")
        print(f"  {os_info['primary_cmd'].format(package=packages)}")
    print("=" * 60 + "\n")
    
    return True


def main():
    # Verify source
    if not os.path.exists(RAW_DIR):
        print(f"Error: Source directory '{RAW_DIR}' not found.")
        return
    
    # Ensure destination directories exist
    for d in [SPACE_DIR, STATE_DIR]:
        if not os.path.exists(d):
            os.makedirs(d)

    # --- Dependency Check ---
    print("[*] Checking dependencies...")
    
    # Check Python modules
    check_python_module("PIL", "PIL", "Pillow")
    check_python_module("pypdf", "pypdf", "pypdf")
    check_python_module("docx", "docx", "python-docx")
    check_python_module("openpyxl", "openpyxl", "openpyxl")
    check_python_module("pptx", "pptx", "python-pptx")
    check_python_module("magic", "magic", "python-magic")
    
    # Check system tools (package names only - OS-specific commands are auto-detected)
    check_system_tool("tesseract", "tesseract-ocr")
    check_system_tool("pdftoppm", "poppler-utils")
    check_system_tool("pdftocairo", "poppler-utils")
    check_system_tool("magick", "imagemagick")
    
    # Print dependency report
    has_missing = print_dependency_report()
    
    # Set up tesseract environment before processing
    tesseract_path = _find_tool("tesseract")
    if tesseract_path:
        print(f"[*] Found tesseract: {tesseract_path}")
        if set_tesseract_env(tesseract_path):
            print(f"[*] TESSDATA_PREFIX set for tesseract")
    else:
        print("[!] Warning: tesseract not found. PDF/image OCR will not work.")

    pdftoppm_path = _find_tool("pdftoppm")
    if pdftoppm_path:
        print(f"[*] Found pdftoppm: {pdftoppm_path}")
    else:
        print("[!] Warning: pdftoppm not found. PDF OCR will not work.")

    processed_files = load_manifest()
    
    to_process = [f for f in os.listdir(RAW_DIR) 
                  if f not in processed_files and not f.startswith('.')]

    if not to_process:
        print("No new files to process.")
        return

    print(f"Found {len(to_process)} candidate files...")
    print("[*] MIME type validation: ENABLED (content-based detection)")
    print()

    results = {"success": [], "mime_warnings": [], "skipped": []}

    for filename in to_process:
        file_path = os.path.join(RAW_DIR, filename)
        
        if os.path.isfile(file_path):
            # Step 1: Detect MIME type from file content
            detected_mime = detect_file_type(file_path)
            ext = os.path.splitext(filename)[1].lower()
            
            # Step 2: Validate MIME type against extension
            is_valid, warning_msg = validate_mime_vs_extension(file_path, detected_mime) if detected_mime else (True, None)
            
            if detected_mime:
                print(f"[*] {filename}: extension='{ext}', detected MIME='{detected_mime}'")
            
            if warning_msg:
                print(f"[WARNING] {filename}: {warning_msg}")
                results["mime_warnings"].append((filename, warning_msg))
            
            # Step 3: Process the file
            if process_file(file_path, filename):
                save_to_manifest(filename)
                print(f"[SUCCESS] Integrated: {filename}")
                results["success"].append(filename)
            else:
                print(f"[SKIPPED] {filename}")
                results["skipped"].append(filename)

    # Print summary
    print()
    print("=" * 60)
    print("INGESTION SUMMARY")
    print("=" * 60)
    print(f"  Files processed successfully: {len(results['success'])}")
    print(f"  Files with MIME warnings:     {len(results['mime_warnings'])}")
    print(f"  Files skipped:                {len(results['skipped'])}")
    if results["mime_warnings"]:
        print()
        print("MIME mismatch details:")
        for fname, msg in results["mime_warnings"]:
            print(f"  - {fname}: {msg}")
    print("=" * 60)

if __name__ == "__main__":
    main()

