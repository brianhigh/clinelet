import os
import re
import sys
import zipfile

# Constants
RAW_DIR = "raw"
SPACE_DIR = "wiki"
STATE_DIR = ".clinelet"
MANIFEST_PATH = os.path.join(STATE_DIR, "processed_files.txt")

# Magic bytes (file signatures) for MIME type detection
MAGIC_SIGNATURES = {
    b'%PDF': 'application/pdf',
    b'PK\x03\x04': 'application/zip',
    b'\x89PNG\r\n\x1a\n': 'image/png',
    b'\xff\xd8\xff': 'image/jpeg',
    b'GIF87a': 'image/gif',
    b'GIF89a': 'image/gif',
    b'BM': 'image/bmp',
    b'II\x2a\x00': 'image/tiff',
    b'MM\x00\x2a': 'image/tiff',
    b'RIFF': 'image/webp',
    b'<!DOCTYPE html': 'text/html',
    b'<html': 'text/html',
    b'<?xml': 'text/xml',
}

# MIME type to extension mapping for validation
MIME_TO_EXTENSIONS = {
    'application/pdf': ['.pdf'],
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': ['.docx'],
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': ['.xlsx'],
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': ['.pptx'],
    'image/png': ['.png'],
    'image/jpeg': ['.jpg', '.jpeg'],
    'image/gif': ['.gif'],
    'image/bmp': ['.bmp'],
    'image/tiff': ['.tiff', '.tif'],
    'image/webp': ['.webp'],
    'text/plain': ['.txt', '.md'],
    'text/html': ['.html', '.htm'],
    'text/xml': ['.xml'],
    'application/zip': ['.zip'],
}

# Try to import python-magic for more robust MIME detection
HAVE_PYTHON_MAGIC = False
try:
    import magic
    HAVE_PYTHON_MAGIC = True
except ImportError:
    pass


def detect_file_type(file_path):
    """Detect the MIME type of a file by examining its magic bytes (file signature).
    
    This provides content-based detection rather than relying on filename extensions,
    which helps catch mismatched or misleading file types.
    
    Args:
        file_path: Path to the file to detect.
        
    Returns:
        A string containing the detected MIME type (e.g., 'application/pdf'),
        or None if the type could not be determined.
    """
    # If python-magic is available, use it for best results
    if HAVE_PYTHON_MAGIC:
        try:
            mime = magic.Magic(mime=True)
            detected = mime.from_file(file_path)
            if detected:
                return detected
        except Exception:
            pass
    
    # Fallback to manual magic byte detection
    try:
        with open(file_path, 'rb') as f:
            header = f.read(12)
            if not header or len(header) < 2:
                return None
            
            # Check fixed signatures first (exact byte match at offset 0)
            for sig, mime_type in MAGIC_SIGNATURES.items():
                if isinstance(sig, bytes) and header[:len(sig)] == sig:
                    if sig == b'RIFF':
                        if len(header) >= 12 and f.read(4) == b'WEBP':
                            return 'image/webp'
                        else:
                            continue
                    return mime_type
            
            # For ZIP-based files (DOCX, XLSX, PPTX), inspect the ZIP contents
            if header[:4] == b'PK\x03\x04':
                try:
                    with zipfile.ZipFile(file_path, 'r') as zf:
                        for name in zf.namelist():
                            name_lower = name.lower()
                            if name_lower.startswith('word/') and name_lower.endswith('.xml'):
                                return 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
                            elif name_lower.startswith('xl/') and name_lower.endswith('.xml'):
                                return 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
                            elif name_lower.startswith('ppt/') and name_lower.endswith('.xml'):
                                return 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
                        return 'application/zip'
                except Exception:
                    return 'application/zip'
            
            # For text files, check the first line
            try:
                first_line = header.decode('utf-8', errors='ignore').strip().lower()
                if first_line.startswith('<!doctype html') or first_line.startswith('<html'):
                    return 'text/html'
                elif first_line.startswith('<?xml'):
                    return 'text/xml'
                elif first_line and all(c.isprintable() or c in '\n\r\t' for c in first_line[:100]):
                    return 'text/plain'
            except Exception:
                pass
            
            return None
    except Exception:
        return None


def validate_mime_vs_extension(file_path, detected_mime):
    """Validate that the file's extension matches its detected MIME type.
    
    Args:
        file_path: Path to the file.
        detected_mime: The detected MIME type string.
        
    Returns:
        A tuple (is_valid, warning_message). is_valid is True if the extension
        matches the MIME type. warning_message describes any mismatch.
    """
    ext = os.path.splitext(file_path)[1].lower()
    
    if detected_mime not in MIME_TO_EXTENSIONS:
        return True, None
    
    expected_extensions = MIME_TO_EXTENSIONS[detected_mime]
    
    if ext in expected_extensions:
        return True, None
    else:
        if detected_mime == 'application/zip':
            return True, None
        if detected_mime == 'text/plain':
            return True, None
        
        return False, f"Extension '{ext}' does not match detected MIME type '{detected_mime}'. Expected one of: {expected_extensions}"


def load_manifest():
    """Loads processed filenames and ensures state directory exists."""
    if not os.path.exists(STATE_DIR):
        os.makedirs(STATE_DIR)
        
    if not os.path.exists(MANIFEST_PATH):
        open(MANIFEST_PATH, 'a').close()
        return set()

    try:
        with open(MANIFEST_PATH, 'r', encoding='utf-8') as f:
            return set(line.strip() for line in f if line.strip())
    except Exception as e:
        print(f"Error reading manifest: {e}")
        return set()


def save_to_manifest(filename):
    """Appends a filename to the manifest file."""
    try:
        if not os.path.exists(STATE_DIR):
            os.makedirs(STATE_DIR)
            
        with open(MANIFEST_PATH, 'a', encoding='utf-8') as f:
            f.write(filename + "\n")
    except Exception as e:
        print(f"Warning: Could not update manifest: {e}")


def to_snake_case(filename):
    """Converts filename to snake_case and ensures it ends in .md"""
    name = os.path.splitext(filename)[0]
    name = re.sub(r'[^a-zA-Z0-9]+', '_', name)
    return name.strip('_').lower() + ".md"


def get_unique_path(target_dir, filename):
    """Prevents overwriting by appending a counter if the file exists."""
    base, ext = os.path.splitext(filename)
    counter = 1
    target_path = os.path.join(target_dir, filename)
    
    while os.path.exists(target_path):
        target_path = os.path.join(target_dir, f"{base}_{counter}{ext}")
        counter += 1
    return target_path


def process_file(file_path, filename):
    """Extracts text and saves to markdown. Returns False if skipped/failed."""
    new_filename = to_snake_case(filename)
    target_path = get_unique_path(SPACE_DIR, new_filename)
    
    ext = os.path.splitext(filename)[1].lower()
    content = ""
    
    # Supported formats for this lite script (non-image, no OCR)
    SUPPORTED_FORMATS = {
        '.md': 'text',
        '.txt': 'text',
        '.html': 'html',
        '.pdf': 'pdf',
        '.docx': 'docx',
        '.xlsx': 'xlsx',
        '.pptx': 'pptx',
    }
    
    if ext not in SUPPORTED_FORMATS:
        print(f"[!] Unsupported format: {ext} ({filename})")
        return False
    
    try:
        if ext in ['.md', '.txt']:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                content = f.read()

        elif ext == '.html':
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                html = f.read()
            content = re.sub(r'<[^>]+>', '', html)

        elif ext == '.pdf':
            try:
                import pypdf
                reader = pypdf.PdfReader(file_path)
                content = "\n".join([page.extract_text() or "" for page in reader.pages]).strip()
                if not content:
                    print(f"[!] Skipped: {filename} (PDF has no text layer/is image-based)")
                    return False
            except ImportError:
                print(f"[-] Missing 'pypdf'. Skipping: {filename}")
                return False

        elif ext == '.docx':
            try:
                import docx
                doc = docx.Document(file_path)
                content = "\n".join([p.text for p in doc.paragraphs]).strip()
            except ImportError:
                print(f"[-] Missing 'python-docx'. Skipping: {filename}")
                return False

        elif ext == '.xlsx':
            try:
                import openpyxl
                wb = openpyxl.load_workbook(file_path, data_only=True)
                lines = []
                for sheet in wb.worksheets:
                    lines.append(f"--- Sheet: {sheet.title} ---")
                    for row in sheet.iter_rows(values_only=True):
                        row_str = " ".join(str(c) for c in row if c is not None).strip()
                        if row_str:
                            lines.append(row_str)
                content = "\n".join(lines)
            except ImportError:
                print(f"[-] Missing 'openpyxl'. Skipping: {filename}")
                return False

        elif ext == '.pptx':
            try:
                import pptx
                prs = pptx.Presentation(file_path)
                lines = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            lines.append(shape.text.strip())
                content = "\n".join(lines)
            except ImportError:
                print(f"[-] Missing 'python-pptx'. Skipping: {filename}")
                return False

    except Exception as e:
        print(f"[ERROR] Critical error processing {filename}: {e}")
        return False
    
    # Final check: If content is still empty after processing
    if not content.strip():
        print(f"[!] Skipped: {filename} (No content extracted)")
        return False

    # Final write to markdown
    with open(target_path, 'w', encoding='utf-8') as f:
        f.write(content)
    
    return True


def main():
    # Verify source
    if not os.path.exists(RAW_DIR):
        print(f"Error: Source directory '{RAW_DIR}' not found.")
        return
    
    # Ensure destination directories exist
    for d in [SPACE_DIR, STATE_DIR]:
        if not os.path.exists(d):
            os.makedirs(d)

    processed_files = load_manifest()
    
    to_process = [f for f in os.listdir(RAW_DIR) 
                  if f not in processed_files and not f.startswith('.')]

    if not to_process:
        print("No new files to process.")
        return

    print(f"Found {len(to_process)} candidate files...")
    print("[*] MIME type validation: ENABLED (content-based detection)")
    print()

    results = {"success": [], "mime_warnings": [], "skipped": []}

    for filename in to_process:
        file_path = os.path.join(RAW_DIR, filename)
        
        if os.path.isfile(file_path):
            # Step 1: Detect MIME type from file content
            detected_mime = detect_file_type(file_path)
            ext = os.path.splitext(filename)[1].lower()
            
            # Step 2: Validate MIME type against extension
            is_valid, warning_msg = validate_mime_vs_extension(file_path, detected_mime) if detected_mime else (True, None)
            
            if detected_mime:
                print(f"[*] {filename}: extension='{ext}', detected MIME='{detected_mime}'")
            
            if warning_msg:
                print(f"[WARNING] {filename}: {warning_msg}")
                results["mime_warnings"].append((filename, warning_msg))
            
            # Step 3: Process the file
            if process_file(file_path, filename):
                save_to_manifest(filename)
                print(f"[SUCCESS] Integrated: {filename}")
                results["success"].append(filename)
            else:
                print(f"[SKIPPED] {filename}")
                results["skipped"].append(filename)

    # Print summary
    print()
    print("=" * 60)
    print("INGESTION SUMMARY")
    print("=" * 60)
    print(f"  Files processed successfully: {len(results['success'])}")
    print(f"  Files with MIME warnings:     {len(results['mime_warnings'])}")
    print(f"  Files skipped:                {len(results['skipped'])}")
    if results["mime_warnings"]:
        print()
        print("MIME mismatch details:")
        for fname, msg in results["mime_warnings"]:
            print(f"  - {fname}: {msg}")
    print("=" * 60)


if __name__ == "__main__":
    main()